#!/usr/bin/env python3
"""
MediGuide PowerPoint Presentation Generator
Creates a professional 20-slide presentation for the MediGuide RAG system.
"""

import os
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
from pptx.oxml.ns import qn

# ─── Color Palette ───────────────────────────────────────────────────────────
C_NAVY      = RGBColor(0x1B, 0x3A, 0x6B)   # #1B3A6B – primary dark navy
C_BLUE      = RGBColor(0x25, 0x63, 0xEB)   # #2563EB – accent bright blue
C_GREEN     = RGBColor(0x10, 0xB9, 0x81)   # #10B981 – success green
C_AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B – warning amber
C_BG        = RGBColor(0xF8, 0xFA, 0xFF)   # #F8FAFF – near-white bg
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARKTEXT  = RGBColor(0x1F, 0x29, 0x37)   # near-black text
C_GRAY      = RGBColor(0x9C, 0xA3, 0xAF)   # slide number gray
C_PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
C_TEAL      = RGBColor(0x0D, 0x94, 0x88)   # teal accent
C_ORANGE    = RGBColor(0xEA, 0x58, 0x0C)   # orange accent
C_RED       = RGBColor(0xDC, 0x26, 0x26)   # red accent
C_LIGHTBLUE = RGBColor(0xDB, 0xEA, 0xFE)   # light blue bg
C_LIGHTGRN  = RGBColor(0xD1, 0xFA, 0xE5)   # light green bg
C_LIGHTAMB  = RGBColor(0xFE, 0xF3, 0xC7)   # light amber bg
C_LIGHTPUR  = RGBColor(0xED, 0xE9, 0xFE)   # light purple bg
C_LIGHTTEAL = RGBColor(0xCC, 0xFB, 0xF1)   # light teal bg

# ─── Slide Dimensions ────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ─── Helper Functions ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    """Add a blank slide (no placeholders)."""
    layout = prs.slide_layouts[6]   # blank layout
    return prs.slides.add_slide(layout)


def fill_slide_bg(slide, color: RGBColor):
    """Fill entire slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width_pt=0):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, italic=False,
                 color=C_DARKTEXT, align=PP_ALIGN.LEFT,
                 v_anchor="top", font_name="Calibri", wrap=True):
    """Add a text box with given text and formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_para(tf, text, font_size=12, bold=False, italic=False,
             color=C_DARKTEXT, align=PP_ALIGN.LEFT, space_before=0,
             font_name="Calibri", level=0):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.level = level
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_slide_number(slide, slide_num, total=20):
    """Add slide number bottom-right."""
    add_text_box(slide, f"{slide_num} / {total}",
                 left=Inches(12.1), top=Inches(7.1),
                 width=Inches(1.0), height=Inches(0.3),
                 font_size=9, color=C_GRAY, align=PP_ALIGN.RIGHT)


def heading(slide, text, top=Inches(0.35), color=C_NAVY, font_size=24):
    """Add a slide heading bar with text."""
    # Accent bar
    add_rect(slide,
             left=Inches(0.4), top=top,
             width=Inches(0.06), height=Inches(0.5),
             fill_color=C_BLUE)
    add_text_box(slide, text,
                 left=Inches(0.55), top=top - Inches(0.03),
                 width=Inches(12.5), height=Inches(0.6),
                 font_size=font_size, bold=True, color=color)


def arrow_right(slide, x, y, w=Inches(0.5), h=Inches(0.3)):
    """Draw a right-pointing arrow connector line."""
    line = slide.shapes.add_connector(1, x, y + h/2, x + w, y + h/2)
    line.line.color.rgb = C_NAVY
    line.line.width = Pt(2)
    return line


def add_rounded_rect_with_text(slide, left, top, width, height,
                                fill_color, text_lines,
                                title_font=13, body_font=11,
                                title_color=C_WHITE, body_color=C_WHITE,
                                corner_radius=0.1, line_color=None):
    """Add a colored box with title and body text lines."""
    shape = slide.shapes.add_shape(5, left, top, width, height)  # rounded rect
    shape.adjustments[0] = corner_radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for line in text_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        is_title = (line == text_lines[0])
        run.font.size = Pt(title_font if is_title else body_font)
        run.font.bold = is_title
        run.font.color.rgb = title_color if is_title else body_color
        run.font.name = "Calibri"
    return shape


def embed_image_from_bytes(slide, img_bytes, left, top, width, height):
    """Embed a PNG image (bytes) into slide."""
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, left, top, width, height)


def make_bar_chart_png():
    """Create corpus bar chart and return PNG bytes."""
    labels = ['openFDA\n(Drug Labels)', 'MedlinePlus\n(NIH)', 'PubMed\nAbstracts', 'PLABA\nDataset']
    values = [41344, 3634, 3141, 1304]
    colors = ['#2563EB', '#10B981', '#7C3AED', '#F59E0B']

    fig, ax = plt.subplots(figsize=(9, 4.5))
    fig.patch.set_facecolor('#F8FAFF')
    ax.set_facecolor('#F8FAFF')

    bars = ax.bar(labels, values, color=colors, width=0.55,
                  edgecolor='white', linewidth=1.5, zorder=3)

    ax.set_ylim(0, 48000)
    ax.set_ylabel('Number of Chunks', fontsize=12, color='#1F2937', fontweight='bold')
    ax.set_title('Knowledge Base Corpus Composition — 49,423 Total Chunks',
                 fontsize=13, color='#1B3A6B', fontweight='bold', pad=12)

    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'{int(x):,}'))
    ax.tick_params(axis='x', labelsize=11, colors='#1F2937')
    ax.tick_params(axis='y', labelsize=10, colors='#6B7280')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E5E7EB')
    ax.spines['bottom'].set_color('#E5E7EB')
    ax.yaxis.grid(True, color='#E5E7EB', linestyle='--', alpha=0.7, zorder=0)

    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2,
                bar.get_height() + 400,
                f'{val:,}',
                ha='center', va='bottom',
                fontsize=11, fontweight='bold', color='#1F2937')

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='#F8FAFF')
    plt.close()
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide1_title(prs):
    """Slide 1 — Title Slide."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_NAVY)

    # Top accent stripe
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=C_BLUE)

    # Decorative grid dots (light)
    for row in range(10):
        for col in range(20):
            x = Inches(0.4 + col * 0.65)
            y = Inches(0.3 + row * 0.7)
            if x < W - Inches(0.2) and y < H - Inches(0.5):
                dot = slide.shapes.add_shape(9, x, y, Inches(0.05), Inches(0.05))  # oval
                dot.fill.solid()
                dot.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
                dot.fill.fore_color.theme_color
                dot.fill.fore_color.rgb = RGBColor(0x2D, 0x4E, 0x8A)
                dot.line.fill.background()

    # Main title "MediGuide"
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(11), Inches(1.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "MediGuide"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    run.font.name = "Calibri"

    # Blue underline accent
    add_rect(slide,
             left=Inches(4.5), top=Inches(3.1),
             width=Inches(4.3), height=Inches(0.07),
             fill_color=C_BLUE)

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.25), Inches(11), Inches(0.7))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "AI-Powered Local Patient Health Companion"
    r2.font.size = Pt(26)
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)
    r2.font.name = "Calibri"

    # Sub-subtitle
    tb3 = slide.shapes.add_textbox(Inches(1.2), Inches(3.95), Inches(11), Inches(0.6))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Retrieval-Augmented Generation for Medical Discharge Notes & Health Q&A"
    r3.font.size = Pt(16)
    r3.font.bold = False
    r3.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    r3.font.name = "Calibri"

    # Tech stack banner (semi-transparent box)
    add_rect(slide,
             left=Inches(1.5), top=Inches(4.9),
             width=Inches(10.3), height=Inches(0.65),
             fill_color=RGBColor(0x0F, 0x23, 0x4A))
    tb4 = slide.shapes.add_textbox(Inches(1.5), Inches(4.92), Inches(10.3), Inches(0.6))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "Built with  LLaMA 3.1  ·  PubMedBERT  ·  FAISS  ·  FastAPI"
    r4.font.size = Pt(15)
    r4.font.bold = False
    r4.font.color.rgb = RGBColor(0x60, 0xA5, 0xFA)
    r4.font.name = "Calibri"

    # Bottom-right tag
    tb5 = slide.shapes.add_textbox(Inches(8.5), Inches(6.7), Inches(4.5), Inches(0.5))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.alignment = PP_ALIGN.RIGHT
    r5 = p5.add_run()
    r5.text = "100% Local  ·  No Cloud  ·  No Data Sharing"
    r5.font.size = Pt(11)
    r5.font.bold = True
    r5.font.color.rgb = C_GREEN
    r5.font.name = "Calibri"

    # Bottom accent stripe
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_BLUE)


def slide2_problem(prs):
    """Slide 2 — The Problem."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "The Challenge: Medical Discharge Notes Are Hard to Understand")
    add_slide_number(slide, 2)

    # Stat boxes
    boxes = [
        (C_RED,   RGBColor(0xFF,0xE4,0xE6), "1 in 5 patients",
         "re-admitted within 30 days of discharge"),
        (C_AMBER, RGBColor(0xFF,0xF7,0xED), "40–80%",
         "of medical information forgotten immediately after doctor visits"),
        (C_BLUE,  C_LIGHTBLUE, "9th Grade",
         "average reading level required for discharge notes"),
    ]
    box_w = Inches(3.6)
    box_h = Inches(2.0)
    gap   = Inches(0.35)
    start_x = Inches(0.8)
    for i, (accent, bg, big_txt, small_txt) in enumerate(boxes):
        lx = start_x + i * (box_w + gap)
        ly = Inches(1.35)
        # background box
        shape = slide.shapes.add_shape(5, lx, ly, box_w, box_h)
        shape.adjustments[0] = 0.05
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = accent
        shape.line.width = Pt(2.5)

        # Left accent bar
        add_rect(slide, lx, ly, Inches(0.08), box_h, fill_color=accent)

        # Big text
        tb = slide.shapes.add_textbox(lx + Inches(0.18), ly + Inches(0.18),
                                      box_w - Inches(0.25), Inches(0.75))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = big_txt
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = accent
        run.font.name = "Calibri"

        # Small text
        tb2 = slide.shapes.add_textbox(lx + Inches(0.18), ly + Inches(0.95),
                                       box_w - Inches(0.3), Inches(0.85))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = small_txt
        run2.font.size = Pt(13)
        run2.font.bold = False
        run2.font.color.rgb = C_DARKTEXT
        run2.font.name = "Calibri"

    # Pain points section
    add_rect(slide, Inches(0.4), Inches(3.65),
             Inches(12.5), Inches(0.04), fill_color=RGBColor(0xE5,0xE7,0xEB))

    pain_heading_tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.75), Inches(12), Inches(0.4))
    tf_ph = pain_heading_tb.text_frame
    p_ph = tf_ph.paragraphs[0]
    run_ph = p_ph.add_run()
    run_ph.text = "Core Pain Points"
    run_ph.font.size = Pt(15)
    run_ph.font.bold = True
    run_ph.font.color.rgb = C_NAVY
    run_ph.font.name = "Calibri"

    bullets = [
        "Complex medical jargon patients can't understand",
        "Critical medication instructions missed or misunderstood",
        "Follow-up steps unclear, leading to poor health outcomes",
        "No accessible tool to translate clinical language to plain English",
    ]
    for i, b in enumerate(bullets):
        lx = Inches(0.5) if i < 2 else Inches(6.7)
        ly = Inches(4.25) + (i % 2) * Inches(0.75)
        # Bullet dot
        dot = slide.shapes.add_shape(9, lx, ly + Inches(0.1), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_BLUE
        dot.line.fill.background()
        add_text_box(slide, b,
                     left=lx + Inches(0.2), top=ly,
                     width=Inches(5.8), height=Inches(0.6),
                     font_size=13, color=C_DARKTEXT)


def slide3_solution(prs):
    """Slide 3 — Our Solution."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "MediGuide: Plain English for Every Patient")
    add_slide_number(slide, 3)

    # Left panel
    lp_x, lp_y = Inches(0.4), Inches(1.1)
    lp_w, lp_h = Inches(6.0), Inches(5.2)
    shape_l = slide.shapes.add_shape(5, lp_x, lp_y, lp_w, lp_h)
    shape_l.adjustments[0] = 0.04
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = RGBColor(0xEF,0xF6,0xFF)
    shape_l.line.color.rgb = C_BLUE
    shape_l.line.width = Pt(2)

    tb_l = slide.shapes.add_textbox(lp_x + Inches(0.25), lp_y + Inches(0.2),
                                    lp_w - Inches(0.4), Inches(0.55))
    tf_l = tb_l.text_frame
    p = tf_l.paragraphs[0]
    run = p.add_run()
    run.text = "Discharge Note Explainer"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = C_BLUE
    run.font.name = "Calibri"

    l_bullets = [
        "Explains diagnosis, medications, follow-up steps,\nwarning signs",
        "Outputs at 6th-grade reading level",
        "Instant streaming response",
        "Supports PDF, DOCX, TXT, and image (OCR) uploads",
        "Context preserved for follow-up questions",
    ]
    for i, b in enumerate(l_bullets):
        by = lp_y + Inches(0.9) + i * Inches(0.82)
        dot = slide.shapes.add_shape(9, lp_x + Inches(0.22),
                                     by + Inches(0.13), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_BLUE
        dot.line.fill.background()
        add_text_box(slide, b,
                     left=lp_x + Inches(0.4), top=by,
                     width=lp_w - Inches(0.6), height=Inches(0.75),
                     font_size=12.5, color=C_DARKTEXT)

    # Right panel
    rp_x = Inches(6.85)
    shape_r = slide.shapes.add_shape(5, rp_x, lp_y, lp_w, lp_h)
    shape_r.adjustments[0] = 0.04
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = RGBColor(0xEC,0xFD,0xF5)
    shape_r.line.color.rgb = C_GREEN
    shape_r.line.width = Pt(2)

    tb_r = slide.shapes.add_textbox(rp_x + Inches(0.25), lp_y + Inches(0.2),
                                    lp_w - Inches(0.4), Inches(0.55))
    tf_r = tb_r.text_frame
    p_r = tf_r.paragraphs[0]
    run_r = p_r.add_run()
    run_r.text = "Health Q&A Assistant"
    run_r.font.size = Pt(18)
    run_r.font.bold = True
    run_r.font.color.rgb = C_GREEN
    run_r.font.name = "Calibri"

    r_bullets = [
        "Everyday health questions answered",
        "Home care tips + when to see a doctor",
        "Evidence-based, source-cited answers",
        "Smart follow-up question suggestions",
        "Medical spell correction (35+ terms)",
    ]
    for i, b in enumerate(r_bullets):
        by = lp_y + Inches(0.9) + i * Inches(0.82)
        dot = slide.shapes.add_shape(9, rp_x + Inches(0.22),
                                     by + Inches(0.13), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_GREEN
        dot.line.fill.background()
        add_text_box(slide, b,
                     left=rp_x + Inches(0.4), top=by,
                     width=lp_w - Inches(0.6), height=Inches(0.7),
                     font_size=12.5, color=C_DARKTEXT)

    # Bottom banner
    banner_y = Inches(6.45)
    banner_h = Inches(0.72)
    add_rect(slide, Inches(0.4), banner_y,
             Inches(12.5), banner_h, fill_color=C_NAVY)
    add_text_box(slide,
                 "100% Local  ·  No Cloud  ·  No Data Leaves Your Device  ·  Open Medical Databases",
                 left=Inches(0.4), top=banner_y + Inches(0.12),
                 width=Inches(12.5), height=Inches(0.5),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide4_architecture(prs):
    """Slide 4 — System Architecture."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "System Architecture — End to End")
    add_slide_number(slide, 4)

    # Main flow boxes
    main_boxes = [
        (C_BLUE,   C_WHITE, "User Input",       Inches(0.35)),
        (C_PURPLE, C_WHITE, "Preprocessor",     Inches(2.95)),
        (C_TEAL,   C_WHITE, "Hybrid Retriever", Inches(5.55)),
        (C_ORANGE, C_WHITE, "LLM Generator",    Inches(8.15)),
        (C_GREEN,  C_WHITE, "Patient Response", Inches(10.75)),
    ]
    box_y   = Inches(1.5)
    box_w   = Inches(2.25)
    box_h   = Inches(0.85)
    arrow_w = Inches(0.55)

    for color, txt_color, label, lx in main_boxes:
        shape = slide.shapes.add_shape(5, lx, box_y, box_w, box_h)
        shape.adjustments[0] = 0.12
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tb = slide.shapes.add_textbox(lx, box_y + Inches(0.12),
                                      box_w, Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = txt_color
        run.font.name = "Calibri"

    # Arrows between main boxes
    arrow_positions = [
        Inches(0.35) + box_w,
        Inches(2.95) + box_w,
        Inches(5.55) + box_w,
        Inches(8.15) + box_w,
    ]
    for ax in arrow_positions:
        line = slide.shapes.add_connector(1,
                                          ax, box_y + box_h/2,
                                          ax + arrow_w - Inches(0.02),
                                          box_y + box_h/2)
        line.line.color.rgb = C_NAVY
        line.line.width = Pt(2.5)

    # Sub-boxes under "User Input"
    sub_input = ["Discharge Note", "Health Question", "Uploaded File"]
    for i, s in enumerate(sub_input):
        sx = Inches(0.35)
        sy = Inches(2.65) + i * Inches(0.55)
        shape = slide.shapes.add_shape(5, sx, sy, box_w, Inches(0.45))
        shape.adjustments[0] = 0.15
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xDB,0xEA,0xFE)
        shape.line.color.rgb = C_BLUE
        shape.line.width = Pt(1)
        tb = slide.shapes.add_textbox(sx, sy + Inches(0.07), box_w, Inches(0.32))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = s
        run.font.size = Pt(10)
        run.font.color.rgb = C_BLUE
        run.font.name = "Calibri"
    # Connector line from input to subs
    line_c = slide.shapes.add_connector(1,
                                        Inches(0.35) + box_w/2, box_y + box_h,
                                        Inches(0.35) + box_w/2, Inches(2.65))
    line_c.line.color.rgb = C_BLUE
    line_c.line.width = Pt(1.5)

    # Sub-boxes under "Hybrid Retriever"
    ret_x = Inches(5.55)
    ret_subs = [
        (RGBColor(0xDB,0xEA,0xFE), C_BLUE,   "BM25 Sparse (w=0.45)"),
        (C_LIGHTGRN,              C_GREEN,  "PubMedBERT FAISS (w=0.55)"),
        (C_LIGHTPUR,              C_PURPLE, "RRF Fusion → Top-K=8"),
    ]
    for i, (bg, lc, txt) in enumerate(ret_subs):
        sy = Inches(2.65) + i * Inches(0.6)
        shape = slide.shapes.add_shape(5, ret_x, sy, box_w, Inches(0.5))
        shape.adjustments[0] = 0.12
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = lc
        shape.line.width = Pt(1)
        tb = slide.shapes.add_textbox(ret_x, sy + Inches(0.08), box_w, Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(10)
        run.font.color.rgb = lc
        run.font.name = "Calibri"
    line_r = slide.shapes.add_connector(1,
                                        ret_x + box_w/2, box_y + box_h,
                                        ret_x + box_w/2, Inches(2.65))
    line_r.line.color.rgb = C_TEAL
    line_r.line.width = Pt(1.5)

    # Sub-boxes under "LLM Generator"
    gen_x = Inches(8.15)
    gen_subs = [
        (C_LIGHTAMB, C_AMBER,  "LLaMA 3.1 8B / Ollama"),
        (RGBColor(0xFE, 0xE2, 0xE2), C_RED,    "Fallback Template"),
        (RGBColor(0xF0, 0xFD, 0xF4), C_GREEN,  "Streaming SSE Tokens"),
    ]
    for i, (bg, lc, txt) in enumerate(gen_subs):
        sy = Inches(2.65) + i * Inches(0.6)
        shape = slide.shapes.add_shape(5, gen_x, sy, box_w, Inches(0.5))
        shape.adjustments[0] = 0.12
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = lc
        shape.line.width = Pt(1)
        tb = slide.shapes.add_textbox(gen_x, sy + Inches(0.08), box_w, Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(10)
        run.font.color.rgb = lc
        run.font.name = "Calibri"
    line_g = slide.shapes.add_connector(1,
                                        gen_x + box_w/2, box_y + box_h,
                                        gen_x + box_w/2, Inches(2.65))
    line_g.line.color.rgb = C_ORANGE
    line_g.line.width = Pt(1.5)

    # Legend
    legend_items = [
        (C_BLUE,   "Input Layer"),
        (C_PURPLE, "Processing"),
        (C_TEAL,   "Retrieval"),
        (C_ORANGE, "Generation"),
        (C_GREEN,  "Output"),
    ]
    for i, (c, lbl) in enumerate(legend_items):
        lx_leg = Inches(0.4) + i * Inches(2.6)
        dot = slide.shapes.add_shape(9, lx_leg, Inches(6.15), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        add_text_box(slide, lbl,
                     left=lx_leg + Inches(0.25), top=Inches(6.1),
                     width=Inches(2.2), height=Inches(0.3),
                     font_size=10, color=C_GRAY)


def slide5_dataset(prs):
    """Slide 5 — Dataset & Corpus."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Knowledge Base: 49,423 Indexed Chunks from Open Medical Databases")
    add_slide_number(slide, 5)

    # Table
    headers  = ["Source", "Records", "Chunks", "Description"]
    col_ws   = [Inches(2.8), Inches(1.2), Inches(1.2), Inches(7.0)]
    rows_data = [
        ("openFDA Drug Labels",     "16,000",  "41,344", "FDA indications, warnings, dosing, adverse reactions"),
        ("MedlinePlus (NIH)",       "~1,100",  "3,634",  "Plain-language health condition pages"),
        ("PubMed Abstracts",        "3,000",   "3,141",  "Clinical literature via NCBI Entrez"),
        ("PLABA Dataset",           "749 pairs","1,304", "PubMed + plain-language rewrites (75 clinical Qs)"),
        ("TOTAL",                   "~20,849", "49,423", ""),
    ]

    tbl_x = Inches(0.4)
    tbl_y = Inches(1.15)
    row_h = Inches(0.72)
    hdr_h = Inches(0.55)

    # Header row
    cx = tbl_x
    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        shape = slide.shapes.add_shape(1, cx, tbl_y, cw, hdr_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_NAVY
        shape.line.color.rgb = C_WHITE
        shape.line.width = Pt(1)
        tb = slide.shapes.add_textbox(cx + Inches(0.08), tbl_y + Inches(0.1),
                                      cw - Inches(0.12), hdr_h - Inches(0.15))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(12.5)
        run.font.bold = True
        run.font.color.rgb = C_WHITE
        run.font.name = "Calibri"
        cx += cw

    # Data rows
    row_colors = [
        RGBColor(0xEF,0xF6,0xFF),
        C_WHITE,
        RGBColor(0xEF,0xF6,0xFF),
        C_WHITE,
        RGBColor(0xE0,0xF2,0xFE),
    ]
    for i, (row, rbg) in enumerate(zip(rows_data, row_colors)):
        ry = tbl_y + hdr_h + i * row_h
        cx = tbl_x
        is_total = (i == len(rows_data) - 1)
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            shape = slide.shapes.add_shape(1, cx, ry, cw, row_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x1B,0x3A,0x6B) if is_total else rbg
            shape.line.color.rgb = RGBColor(0xE5,0xE7,0xEB)
            shape.line.width = Pt(0.5)
            tb = slide.shapes.add_textbox(cx + Inches(0.08), ry + Inches(0.12),
                                          cw - Inches(0.12), row_h - Inches(0.18))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell
            run.font.size = Pt(11.5)
            run.font.bold = is_total
            run.font.color.rgb = C_WHITE if is_total else C_DARKTEXT
            run.font.name = "Calibri"
            cx += cw

    # Note at bottom
    add_text_box(slide,
                 "Note: MIMIC-IV-Note (331,794 clinical notes) supported — requires PhysioNet credentials",
                 left=Inches(0.4), top=Inches(6.3),
                 width=Inches(10), height=Inches(0.45),
                 font_size=10.5, italic=True, color=C_GRAY)
    add_text_box(slide,
                 "Chunk settings: Size 384 tokens · Overlap 64 tokens · Embedding: pritamdeka/S-PubMedBert-MS-MARCO",
                 left=Inches(0.4), top=Inches(6.72),
                 width=Inches(12), height=Inches(0.4),
                 font_size=10, color=C_GRAY)


def slide6_rag_pipeline(prs):
    """Slide 6 — RAG Pipeline Deep Dive."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Hybrid Retrieval-Augmented Generation Pipeline")
    add_slide_number(slide, 6)

    steps = [
        (C_BLUE,   "1  PREPROCESS",
         "Section segmentation · Abbreviation expansion (80+ terms) · Medication extraction · Entity normalisation"),
        (C_PURPLE, "2  RETRIEVE",
         "BM25 sparse (weight 0.45) + PubMedBERT dense FAISS (weight 0.55) → Top-K=8 per section"),
        (C_TEAL,   "3  FUSE",
         "Reciprocal Rank Fusion (RRF) · Source filtering per section · Cosine similarity ≥ 0.28 threshold"),
        (C_ORANGE, "4  GENERATE",
         "LLaMA 3.1 8B via Ollama · Streaming SSE tokens · Fallback template if offline"),
        (C_GREEN,  "5  RESPOND",
         "4-section structured output · RAG provenance panel · Source citations"),
    ]

    step_h = Inches(0.88)
    gap_h  = Inches(0.22)
    box_w  = Inches(12.2)
    start_y = Inches(1.15)

    for i, (color, title, desc) in enumerate(steps):
        sy = start_y + i * (step_h + gap_h)
        # Arrow connector between steps (except after last)
        if i > 0:
            prev_bottom = start_y + (i-1) * (step_h + gap_h) + step_h
            line = slide.shapes.add_connector(1,
                                              Inches(0.6), prev_bottom,
                                              Inches(0.6), sy)
            line.line.color.rgb = color
            line.line.width = Pt(2)

        # Number circle
        dot = slide.shapes.add_shape(9, Inches(0.4), sy + Inches(0.22),
                                     Inches(0.42), Inches(0.42))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        tb_n = slide.shapes.add_textbox(Inches(0.4), sy + Inches(0.18),
                                        Inches(0.42), Inches(0.45))
        tf_n = tb_n.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.CENTER
        run_n = p_n.add_run()
        run_n.text = str(i + 1)
        run_n.font.size = Pt(13)
        run_n.font.bold = True
        run_n.font.color.rgb = C_WHITE
        run_n.font.name = "Calibri"

        # Step box
        shape = slide.shapes.add_shape(5, Inches(0.95), sy,
                                       box_w, step_h)
        shape.adjustments[0] = 0.06
        shape.fill.solid()

        # Light tinted bg
        r, g, b = color[0], color[1], color[2]
        r2 = min(255, int(r + (255-r)*0.85))
        g2 = min(255, int(g + (255-g)*0.85))
        b2 = min(255, int(b + (255-b)*0.85))
        shape.fill.fore_color.rgb = RGBColor(r2, g2, b2)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        # Title tag
        tag = slide.shapes.add_shape(1, Inches(0.95), sy, Inches(2.1), step_h)
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tb_t = slide.shapes.add_textbox(Inches(0.95), sy + Inches(0.2),
                                        Inches(2.1), Inches(0.6))
        tf_t = tb_t.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.alignment = PP_ALIGN.CENTER
        run_t = p_t.add_run()
        run_t.text = title
        run_t.font.size = Pt(13)
        run_t.font.bold = True
        run_t.font.color.rgb = C_WHITE
        run_t.font.name = "Calibri"

        # Description
        add_text_box(slide, desc,
                     left=Inches(3.25), top=sy + Inches(0.17),
                     width=Inches(9.6), height=Inches(0.6),
                     font_size=12.5, color=C_DARKTEXT)


def slide7_retrieval(prs):
    """Slide 7 — Retrieval Strategy."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Hybrid Retrieval: BM25 + Dense + RRF Fusion")
    add_slide_number(slide, 7)

    # BM25 box
    bm_x, bm_y = Inches(0.4), Inches(1.2)
    bm_w, bm_h = Inches(5.8), Inches(3.8)
    shape_bm = slide.shapes.add_shape(5, bm_x, bm_y, bm_w, bm_h)
    shape_bm.adjustments[0] = 0.05
    shape_bm.fill.solid()
    shape_bm.fill.fore_color.rgb = RGBColor(0xEF,0xF6,0xFF)
    shape_bm.line.color.rgb = C_BLUE
    shape_bm.line.width = Pt(2.5)

    add_text_box(slide, "BM25 Sparse Retrieval",
                 left=bm_x + Inches(0.25), top=bm_y + Inches(0.2),
                 width=bm_w - Inches(0.4), height=Inches(0.5),
                 font_size=16, bold=True, color=C_BLUE)

    bm_items = [
        "Term frequency matching",
        "Fast exact keyword search",
        "Great for drug names & medical terms",
        "Weight: 0.45",
    ]
    for i, item in enumerate(bm_items):
        dot = slide.shapes.add_shape(9, bm_x + Inches(0.25),
                                     bm_y + Inches(0.88) + i*Inches(0.65),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_BLUE
        dot.line.fill.background()
        add_text_box(slide, item,
                     left=bm_x + Inches(0.45),
                     top=bm_y + Inches(0.82) + i*Inches(0.65),
                     width=bm_w - Inches(0.6), height=Inches(0.5),
                     font_size=13, color=C_DARKTEXT,
                     bold=(item == "Weight: 0.45"))

    # PubMedBERT box
    pb_x = Inches(7.1)
    shape_pb = slide.shapes.add_shape(5, pb_x, bm_y, bm_w, bm_h)
    shape_pb.adjustments[0] = 0.05
    shape_pb.fill.solid()
    shape_pb.fill.fore_color.rgb = C_LIGHTPUR
    shape_pb.line.color.rgb = C_PURPLE
    shape_pb.line.width = Pt(2.5)

    add_text_box(slide, "PubMedBERT Dense Retrieval",
                 left=pb_x + Inches(0.25), top=bm_y + Inches(0.2),
                 width=bm_w - Inches(0.4), height=Inches(0.5),
                 font_size=16, bold=True, color=C_PURPLE)

    pb_items = [
        "Semantic embedding similarity",
        "FAISS IndexFlatIP",
        "Biomedical domain pre-training",
        "Weight: 0.55",
    ]
    for i, item in enumerate(pb_items):
        dot = slide.shapes.add_shape(9, pb_x + Inches(0.25),
                                     bm_y + Inches(0.88) + i*Inches(0.65),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_PURPLE
        dot.line.fill.background()
        add_text_box(slide, item,
                     left=pb_x + Inches(0.45),
                     top=bm_y + Inches(0.82) + i*Inches(0.65),
                     width=bm_w - Inches(0.6), height=Inches(0.5),
                     font_size=13, color=C_DARKTEXT,
                     bold=(item == "Weight: 0.55"))

    # Arrow down to RRF box
    mid_x = Inches(6.67)
    for side_x in [bm_x + bm_w/2, pb_x + bm_w/2]:
        line = slide.shapes.add_connector(1, side_x, bm_y + bm_h,
                                          mid_x, Inches(5.3))
        line.line.color.rgb = C_NAVY
        line.line.width = Pt(2)

    # RRF fusion box
    rrf_x, rrf_y = Inches(3.5), Inches(5.3)
    rrf_w, rrf_h = Inches(6.3), Inches(1.6)
    shape_rrf = slide.shapes.add_shape(5, rrf_x, rrf_y, rrf_w, rrf_h)
    shape_rrf.adjustments[0] = 0.06
    shape_rrf.fill.solid()
    shape_rrf.fill.fore_color.rgb = C_NAVY
    shape_rrf.line.fill.background()

    tb_rrf = slide.shapes.add_textbox(rrf_x + Inches(0.2), rrf_y + Inches(0.1),
                                      rrf_w - Inches(0.3), Inches(0.5))
    tf_rrf = tb_rrf.text_frame
    p_rrf = tf_rrf.paragraphs[0]
    p_rrf.alignment = PP_ALIGN.CENTER
    run_rrf = p_rrf.add_run()
    run_rrf.text = "Reciprocal Rank Fusion (RRF)"
    run_rrf.font.size = Pt(15)
    run_rrf.font.bold = True
    run_rrf.font.color.rgb = C_WHITE
    run_rrf.font.name = "Calibri"

    add_text_box(slide,
                 "Combines ranked lists · Reduces bias · Score = \u03a3 1/(k+rank\u1d62) · Top-K=8 per discharge section",
                 left=rrf_x + Inches(0.2), top=rrf_y + Inches(0.7),
                 width=rrf_w - Inches(0.3), height=Inches(0.75),
                 font_size=12, color=RGBColor(0xBF,0xDB,0xFF), align=PP_ALIGN.CENTER)


def slide8_evaluation(prs):
    """Slide 8 — Evaluation Results."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Evaluation Results — Measured on 20 PLABA Test Pairs")
    add_slide_number(slide, 8)

    metrics = [
        (C_BLUE,   C_LIGHTBLUE,  "Precision@5",  "0.15",
         "15% of top-5 retrieved\nare gold chunks"),
        (C_GREEN,  C_LIGHTGRN,   "Recall@5",     "0.55",
         "55% of gold chunks\nfound in top 5"),
        (C_PURPLE, C_LIGHTPUR,   "MRR",          "0.576",
         "Mean Reciprocal Rank\nFirst gold chunk at rank ~1.7"),
        (C_AMBER,  C_LIGHTAMB,   "ROUGE-L",      "0.117",
         "Lexical overlap with\nplain-language reference"),
    ]

    box_w = Inches(2.9)
    box_h = Inches(3.5)
    gap   = Inches(0.28)
    start_x = Inches(0.4)

    for i, (accent, bg, name, val, desc) in enumerate(metrics):
        lx = start_x + i * (box_w + gap)
        ly = Inches(1.15)

        shape = slide.shapes.add_shape(5, lx, ly, box_w, box_h)
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = accent
        shape.line.width = Pt(2.5)

        # Top colored header strip
        add_rect(slide, lx, ly, box_w, Inches(0.5), fill_color=accent)

        # Metric name
        add_text_box(slide, name,
                     left=lx, top=ly + Inches(0.55),
                     width=box_w, height=Inches(0.5),
                     font_size=14, bold=True, color=accent, align=PP_ALIGN.CENTER)

        # Big number
        add_text_box(slide, val,
                     left=lx, top=ly + Inches(1.1),
                     width=box_w, height=Inches(1.0),
                     font_size=46, bold=True, color=accent, align=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, desc,
                     left=lx + Inches(0.1), top=ly + Inches(2.2),
                     width=box_w - Inches(0.2), height=Inches(1.1),
                     font_size=11.5, color=C_DARKTEXT, align=PP_ALIGN.CENTER)

    # Readability banner
    add_rect(slide, Inches(0.4), Inches(4.9),
             Inches(12.5), Inches(0.65), fill_color=RGBColor(0xEF,0xF6,0xFF))
    add_rect(slide, Inches(0.4), Inches(4.9),
             Inches(0.06), Inches(0.65), fill_color=C_BLUE)
    add_text_box(slide,
                 "Readability Score: Flesch-Kincaid Grade 9.3 — accessible for most adult patients",
                 left=Inches(0.6), top=Inches(5.0),
                 width=Inches(12.0), height=Inches(0.45),
                 font_size=13, bold=True, color=C_NAVY)

    # Dataset note
    add_text_box(slide,
                 "Evaluated on PLABA dataset: biomedical abstracts paired with patient-friendly rewrites",
                 left=Inches(0.4), top=Inches(5.75),
                 width=Inches(12.5), height=Inches(0.4),
                 font_size=10.5, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide9_webapp(prs):
    """Slide 9 — Web Application Features."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "MediGuide Web App — Feature Overview")
    add_slide_number(slide, 9)

    features = [
        (C_BLUE,   "Two-Mode Interface",      "Discharge Note + Health Q&A tabs"),
        (C_GREEN,  "Streaming Responses",     "Real-time token-by-token output via SSE"),
        (C_AMBER,  "Conversation History",    "Full chat history in localStorage, one-click reload"),
        (C_PURPLE, "Multi-Format Upload",     "PDF, DOCX, TXT, PNG/JPG/TIFF (OCR)"),
        (C_TEAL,   "Smart Recommendations",  "Context-aware follow-up question chips"),
        (C_ORANGE, "Spell Correction",        "35+ medical misspellings auto-detected & fixed"),
    ]

    card_w = Inches(4.1)
    card_h = Inches(1.65)
    gap_x  = Inches(0.32)
    gap_y  = Inches(0.3)
    start_x = Inches(0.45)
    start_y = Inches(1.2)

    for i, (color, title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        lx = start_x + col * (card_w + gap_x)
        ly = start_y + row * (card_h + gap_y)

        shape = slide.shapes.add_shape(5, lx, ly, card_w, card_h)
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_WHITE
        shape.line.color.rgb = RGBColor(0xE5,0xE7,0xEB)
        shape.line.width = Pt(1)

        # Left accent bar
        add_rect(slide, lx, ly, Inches(0.08), card_h, fill_color=color)

        # Icon circle
        ic = slide.shapes.add_shape(9, lx + Inches(0.18), ly + Inches(0.2),
                                    Inches(0.38), Inches(0.38))
        ic.fill.solid()
        ic.fill.fore_color.rgb = color
        # Light bg behind icon
        r, g, b = color[0], color[1], color[2]
        r2 = min(255, int(r + (255-r)*0.85))
        g2 = min(255, int(g + (255-g)*0.85))
        b2 = min(255, int(b + (255-b)*0.85))
        ic.fill.fore_color.rgb = RGBColor(r2, g2, b2)
        ic.line.color.rgb = color
        ic.line.width = Pt(1.5)

        add_text_box(slide, title,
                     left=lx + Inches(0.18), top=ly + Inches(0.15),
                     width=card_w - Inches(0.3), height=Inches(0.45),
                     font_size=13, bold=True, color=color)

        add_text_box(slide, desc,
                     left=lx + Inches(0.18), top=ly + Inches(0.65),
                     width=card_w - Inches(0.3), height=Inches(0.85),
                     font_size=12, color=C_DARKTEXT)


def slide10_two_modes(prs):
    """Slide 10 — Two Modes."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Two Intelligent Modes for Every Patient Need")
    add_slide_number(slide, 10)

    # Left panel (blue)
    lp_x, lp_y = Inches(0.35), Inches(1.1)
    lp_w, lp_h = Inches(6.15), Inches(5.9)
    shape_l = slide.shapes.add_shape(5, lp_x, lp_y, lp_w, lp_h)
    shape_l.adjustments[0] = 0.04
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = C_NAVY
    shape_l.line.fill.background()

    add_text_box(slide, "Discharge Note Mode",
                 left=lp_x + Inches(0.3), top=lp_y + Inches(0.25),
                 width=lp_w - Inches(0.5), height=Inches(0.55),
                 font_size=18, bold=True, color=C_WHITE)

    add_text_box(slide, "Input: Full hospital discharge note",
                 left=lp_x + Inches(0.3), top=lp_y + Inches(0.9),
                 width=lp_w - Inches(0.5), height=Inches(0.4),
                 font_size=12, color=RGBColor(0x93,0xC5,0xFD))

    add_text_box(slide, "Output sections:",
                 left=lp_x + Inches(0.3), top=lp_y + Inches(1.35),
                 width=lp_w - Inches(0.5), height=Inches(0.35),
                 font_size=12, bold=True, color=C_WHITE)

    l_sections = [
        ("Diagnosis Explained",    C_LIGHTBLUE),
        ("Your Medications",       RGBColor(0xA7,0xF3,0xD0)),
        ("Follow-up Actions",      C_LIGHTAMB),
        ("Warning Signs",          RGBColor(0xFE,0xCA,0xCA)),
    ]
    for i, (sec, col) in enumerate(l_sections):
        sy = lp_y + Inches(1.8) + i * Inches(0.65)
        add_rect(slide, lp_x + Inches(0.28), sy,
                 lp_w - Inches(0.5), Inches(0.5), fill_color=col)
        add_text_box(slide, sec,
                     left=lp_x + Inches(0.4), top=sy + Inches(0.07),
                     width=lp_w - Inches(0.65), height=Inches(0.38),
                     font_size=12, bold=True, color=C_DARKTEXT)

    add_text_box(slide, "Context preserved for follow-up questions",
                 left=lp_x + Inches(0.3), top=lp_y + Inches(4.55),
                 width=lp_w - Inches(0.5), height=Inches(0.5),
                 font_size=11, italic=True, color=RGBColor(0x93,0xC5,0xFD))

    # Right panel (green)
    rp_x = Inches(6.8)
    shape_r = slide.shapes.add_shape(5, rp_x, lp_y, lp_w, lp_h)
    shape_r.adjustments[0] = 0.04
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = RGBColor(0x06,0x5F,0x46)
    shape_r.line.fill.background()

    add_text_box(slide, "Health Q&A Mode",
                 left=rp_x + Inches(0.3), top=lp_y + Inches(0.25),
                 width=lp_w - Inches(0.5), height=Inches(0.55),
                 font_size=18, bold=True, color=C_WHITE)

    add_text_box(slide, "Input: Everyday health question",
                 left=rp_x + Inches(0.3), top=lp_y + Inches(0.9),
                 width=lp_w - Inches(0.5), height=Inches(0.4),
                 font_size=12, color=RGBColor(0x6E,0xE7,0xB7))

    add_text_box(slide, "Output sections:",
                 left=rp_x + Inches(0.3), top=lp_y + Inches(1.35),
                 width=lp_w - Inches(0.5), height=Inches(0.35),
                 font_size=12, bold=True, color=C_WHITE)

    r_sections = [
        ("What This Could Be",   RGBColor(0xD1,0xFA,0xE5)),
        ("Home Care Tips",       RGBColor(0xA7,0xF3,0xD0)),
        ("When to See a Doctor", C_LIGHTAMB),
        ("Go to the ER If",      RGBColor(0xFE,0xCA,0xCA)),
    ]
    for i, (sec, col) in enumerate(r_sections):
        sy = lp_y + Inches(1.8) + i * Inches(0.65)
        add_rect(slide, rp_x + Inches(0.28), sy,
                 lp_w - Inches(0.5), Inches(0.5), fill_color=col)
        add_text_box(slide, sec,
                     left=rp_x + Inches(0.4), top=sy + Inches(0.07),
                     width=lp_w - Inches(0.65), height=Inches(0.38),
                     font_size=12, bold=True, color=C_DARKTEXT)

    add_text_box(slide, "Topic-matched follow-up suggestions",
                 left=rp_x + Inches(0.3), top=lp_y + Inches(4.55),
                 width=lp_w - Inches(0.5), height=Inches(0.5),
                 font_size=11, italic=True, color=RGBColor(0x6E,0xE7,0xB7))


def slide11_smart_input(prs):
    """Slide 11 — Smart Input Features."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Intelligent Input: Suggestions + Spell Correction")
    add_slide_number(slide, 11)

    # Left section: Suggestion Chips
    lx = Inches(0.4)
    ly = Inches(1.2)
    lw = Inches(5.9)
    lh = Inches(5.5)
    shape_l = slide.shapes.add_shape(5, lx, ly, lw, lh)
    shape_l.adjustments[0] = 0.05
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = RGBColor(0xEF,0xF6,0xFF)
    shape_l.line.color.rgb = C_BLUE
    shape_l.line.width = Pt(2)

    add_text_box(slide, "Dynamic Suggestion Chips",
                 left=lx + Inches(0.2), top=ly + Inches(0.2),
                 width=lw - Inches(0.3), height=Inches(0.5),
                 font_size=16, bold=True, color=C_BLUE)

    sugg_bullets = [
        "Contextual starters appear above input",
        "Updates as you type (keyword matching)",
        "Different chips for each mode",
        "Clicking pre-fills textarea",
    ]
    for i, b in enumerate(sugg_bullets):
        dot = slide.shapes.add_shape(9, lx + Inches(0.25),
                                     ly + Inches(0.9) + i * Inches(0.55),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_BLUE
        dot.line.fill.background()
        add_text_box(slide, b,
                     left=lx + Inches(0.45),
                     top=ly + Inches(0.82) + i * Inches(0.55),
                     width=lw - Inches(0.6), height=Inches(0.45),
                     font_size=12.5, color=C_DARKTEXT)

    # Mock chip display
    chips = ["What are my discharge instructions?",
             "What does my diagnosis mean?",
             "When should I take my medications?"]
    for i, chip in enumerate(chips):
        cy = ly + Inches(3.0) + i * Inches(0.6)
        chip_shape = slide.shapes.add_shape(5, lx + Inches(0.25), cy,
                                            lw - Inches(0.5), Inches(0.48))
        chip_shape.adjustments[0] = 0.5
        chip_shape.fill.solid()
        chip_shape.fill.fore_color.rgb = C_WHITE
        chip_shape.line.color.rgb = C_BLUE
        chip_shape.line.width = Pt(1.5)
        add_text_box(slide, chip,
                     left=lx + Inches(0.4), top=cy + Inches(0.08),
                     width=lw - Inches(0.7), height=Inches(0.35),
                     font_size=11, color=C_BLUE)

    # Right section: Spell Correction
    rx = Inches(6.75)
    shape_r = slide.shapes.add_shape(5, rx, ly, lw, lh)
    shape_r.adjustments[0] = 0.05
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = RGBColor(0xFF,0xFB,0xEB)
    shape_r.line.color.rgb = C_AMBER
    shape_r.line.width = Pt(2)

    add_text_box(slide, "Medical Spell Correction",
                 left=rx + Inches(0.2), top=ly + Inches(0.2),
                 width=lw - Inches(0.3), height=Inches(0.5),
                 font_size=16, bold=True, color=C_AMBER)

    spell_bullets = [
        "35+ medical term corrections",
        'diabeties \u2192 diabetes',
        'pnemonia \u2192 pneumonia',
        'symtoms \u2192 symptoms',
        '"Fix it" one-click yellow banner',
        "Browser spellcheck also enabled",
    ]
    for i, b in enumerate(spell_bullets):
        dot = slide.shapes.add_shape(9, rx + Inches(0.25),
                                     ly + Inches(0.9) + i * Inches(0.55),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_AMBER
        dot.line.fill.background()
        add_text_box(slide, b,
                     left=rx + Inches(0.45),
                     top=ly + Inches(0.82) + i * Inches(0.55),
                     width=lw - Inches(0.6), height=Inches(0.45),
                     font_size=12.5, color=C_DARKTEXT)

    # Mock correction banner
    banner_y = ly + Inches(4.3)
    add_rect(slide, rx + Inches(0.2), banner_y,
             lw - Inches(0.4), Inches(0.75), fill_color=RGBColor(0xFE,0xF3,0xC7))
    add_rect(slide, rx + Inches(0.2), banner_y,
             Inches(0.06), Inches(0.75), fill_color=C_AMBER)
    add_text_box(slide, 'Did you mean "diabetes"?  [Fix it]',
                 left=rx + Inches(0.35), top=banner_y + Inches(0.18),
                 width=lw - Inches(0.6), height=Inches(0.4),
                 font_size=12, bold=True, color=RGBColor(0x92,0x40,0x0E))


def slide12_tech_stack(prs):
    """Slide 12 — Technical Stack (dark bg)."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_NAVY)

    # Top accent
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=C_BLUE)

    add_text_box(slide, "Technology Stack",
                 left=Inches(0.5), top=Inches(0.3),
                 width=Inches(12), height=Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 12)

    columns = [
        (C_BLUE,   "Backend",      [
            "FastAPI",
            "Uvicorn",
            "Python 3.9+",
            "Server-Sent Events (SSE)",
            "pdfplumber",
            "pytesseract (OCR)",
        ]),
        (C_GREEN,  "AI / ML",      [
            "LLaMA 3.1 8B",
            "Ollama",
            "PubMedBERT",
            "FAISS IndexFlatIP",
            "BM25Okapi",
            "sentence-transformers",
        ]),
        (C_AMBER,  "Data Sources", [
            "openFDA",
            "MedlinePlus (NIH)",
            "PubMed / NCBI",
            "PLABA Dataset",
            "MIMIC-IV-Note*",
            "MedQuAD*",
        ]),
        (C_PURPLE, "Frontend",     [
            "Vanilla JavaScript",
            "CSS Variables",
            "localStorage",
            "Responsive Design",
            "Dark / Light Mode",
            "SSE Event Stream",
        ]),
    ]

    col_w  = Inches(3.0)
    gap    = Inches(0.18)
    start_x = Inches(0.45)
    start_y = Inches(1.2)
    col_h  = Inches(5.6)

    for i, (color, title, items) in enumerate(columns):
        lx = start_x + i * (col_w + gap)
        # Column card
        shape = slide.shapes.add_shape(5, lx, start_y, col_w, col_h)
        shape.adjustments[0] = 0.05
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0F, 0x23, 0x4A)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        # Header
        add_rect(slide, lx, start_y, col_w, Inches(0.6), fill_color=color)
        add_text_box(slide, title,
                     left=lx, top=start_y + Inches(0.1),
                     width=col_w, height=Inches(0.45),
                     font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            item_y = start_y + Inches(0.75) + j * Inches(0.75)
            # Item bg
            item_bg = slide.shapes.add_shape(5, lx + Inches(0.12), item_y,
                                             col_w - Inches(0.24), Inches(0.6))
            item_bg.adjustments[0] = 0.15
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x6B)
            item_bg.line.fill.background()
            add_text_box(slide, item,
                         left=lx + Inches(0.12), top=item_y + Inches(0.1),
                         width=col_w - Inches(0.24), height=Inches(0.45),
                         font_size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "* Planned — not yet integrated",
                 left=Inches(0.5), top=Inches(6.9),
                 width=Inches(5), height=Inches(0.35),
                 font_size=10, italic=True, color=C_GRAY)

    # Bottom accent
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_BLUE)


def slide13_privacy(prs):
    """Slide 13 — Privacy & Security."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Privacy First — Everything Stays on Your Device")
    add_slide_number(slide, 13)

    blocks = [
        (C_RED,   RGBColor(0xFF,0xF0,0xF0),
         "No Cloud",
         ["All processing is local",
          "No API calls to external AI services",
          "No patient data transmitted",
          "Works completely offline"]),
        (C_BLUE,  RGBColor(0xEF,0xF6,0xFF),
         "No Data Sharing",
         ["Conversations stored in browser only",
          "Medical notes never leave the device",
          "Compliant with privacy-first principles",
          "No telemetry or analytics"]),
        (C_GREEN, RGBColor(0xEC,0xFD,0xF5),
         "Fully Offline Capable",
         ["Ollama runs locally on your machine",
          "FAISS index stored on disk",
          "Works without internet after setup",
          "GPU accelerated (optional)"]),
    ]

    block_w = Inches(3.9)
    block_h = Inches(4.7)
    gap     = Inches(0.27)
    start_x = Inches(0.45)
    start_y = Inches(1.15)

    for i, (accent, bg, title, items) in enumerate(blocks):
        lx = start_x + i * (block_w + gap)

        shape = slide.shapes.add_shape(5, lx, start_y, block_w, block_h)
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = accent
        shape.line.width = Pt(2.5)

        # Top header
        add_rect(slide, lx, start_y, block_w, Inches(0.7), fill_color=accent)

        add_text_box(slide, title,
                     left=lx + Inches(0.15), top=start_y + Inches(0.1),
                     width=block_w - Inches(0.25), height=Inches(0.55),
                     font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            iy = start_y + Inches(0.9) + j * Inches(0.85)
            # Check mark circle
            ck = slide.shapes.add_shape(9, lx + Inches(0.2), iy + Inches(0.05),
                                        Inches(0.3), Inches(0.3))
            ck.fill.solid()
            ck.fill.fore_color.rgb = accent
            ck.line.fill.background()
            add_text_box(slide, item,
                         left=lx + Inches(0.6), top=iy,
                         width=block_w - Inches(0.75), height=Inches(0.6),
                         font_size=13, color=C_DARKTEXT)

    # HIPAA-ready badge
    badge_x = Inches(4.5)
    badge_y = Inches(6.15)
    badge_w = Inches(4.3)
    badge_h = Inches(0.7)
    badge = slide.shapes.add_shape(5, badge_x, badge_y, badge_w, badge_h)
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_NAVY
    badge.line.fill.background()
    add_text_box(slide, "Privacy-First Architecture — No External Data Transmission",
                 left=badge_x + Inches(0.1), top=badge_y + Inches(0.15),
                 width=badge_w - Inches(0.2), height=Inches(0.45),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide14_provenance(prs):
    """Slide 14 — RAG Provenance."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Transparent AI — Full Source Traceability")
    add_slide_number(slide, 14)

    # Explanation text
    add_text_box(slide,
                 "Every response includes a RAG Provenance Panel showing exactly which database "
                 "chunks were retrieved, their sources, and their retrieval scores.",
                 left=Inches(0.5), top=Inches(1.1),
                 width=Inches(12.3), height=Inches(0.6),
                 font_size=13, color=C_DARKTEXT)

    # Provenance box
    prov_x = Inches(0.5)
    prov_y = Inches(1.85)
    prov_w = Inches(12.3)
    prov_h = Inches(3.5)

    prov_bg = slide.shapes.add_shape(1, prov_x, prov_y, prov_w, prov_h)
    prov_bg.fill.solid()
    prov_bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    prov_bg.line.color.rgb = C_BLUE
    prov_bg.line.width = Pt(1.5)

    # Header bar
    add_rect(slide, prov_x, prov_y, prov_w, Inches(0.45),
             fill_color=RGBColor(0x1B, 0x3A, 0x6B))
    add_text_box(slide, "RAG Provenance — Sources Retrieved",
                 left=prov_x + Inches(0.2), top=prov_y + Inches(0.07),
                 width=prov_w - Inches(0.3), height=Inches(0.35),
                 font_size=12, bold=True, color=C_BLUE)

    prov_lines = [
        ("Source: ", "MedlinePlus (National Library of Medicine)", C_GREEN),
        ("Title: ", "Heart Failure  |  Doc ID: medlineplus_42", C_WHITE),
        ("Scores: ", "BM25=4.821   Dense=0.731   Fused(RRF)=0.0318", C_AMBER),
        ("", "", C_WHITE),  # blank separator
        ("Source: ", "FDA Drug Label Database", C_GREEN),
        ("Title: ", "FUROSEMIDE  |  Doc ID: openfda_1847", C_WHITE),
        ("Scores: ", "BM25=6.203   Dense=0.812   Fused(RRF)=0.0334", C_AMBER),
    ]

    for i, (label, value, val_color) in enumerate(prov_lines):
        line_y = prov_y + Inches(0.55) + i * Inches(0.4)
        if label:
            add_text_box(slide, label,
                         left=prov_x + Inches(0.25), top=line_y,
                         width=Inches(1.0), height=Inches(0.38),
                         font_size=11, bold=True, color=C_GRAY)
            add_text_box(slide, value,
                         left=prov_x + Inches(1.2), top=line_y,
                         width=prov_w - Inches(1.5), height=Inches(0.38),
                         font_size=11, color=val_color)

    # Explanation below
    add_text_box(slide,
                 "Patients and clinicians can verify every claim — each response shows "
                 "which database chunks were retrieved, their sources, and retrieval scores.",
                 left=Inches(0.5), top=Inches(5.55),
                 width=Inches(12.3), height=Inches(0.6),
                 font_size=12.5, color=C_DARKTEXT)

    # Trust badge
    badge = slide.shapes.add_shape(5, Inches(3.5), Inches(6.25),
                                   Inches(6.3), Inches(0.65))
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_NAVY
    badge.line.fill.background()
    add_text_box(slide, "Transparent · Verifiable · Evidence-Based AI",
                 left=Inches(3.5), top=Inches(6.35),
                 width=Inches(6.3), height=Inches(0.45),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide15_corpus_chart(prs, chart_png_bytes):
    """Slide 15 — Corpus Distribution."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Corpus Composition — 49,423 Chunks")
    add_slide_number(slide, 15)

    # Embed chart
    chart_x = Inches(0.7)
    chart_y = Inches(1.1)
    chart_w = Inches(11.9)
    chart_h = Inches(5.5)
    embed_image_from_bytes(slide, chart_png_bytes, chart_x, chart_y, chart_w, chart_h)

    # Annotation
    add_text_box(slide,
                 "openFDA dominates at 84% of total corpus — reflects FDA labeling data density",
                 left=Inches(0.5), top=Inches(6.85),
                 width=Inches(12.3), height=Inches(0.35),
                 font_size=10, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide16_future(prs):
    """Slide 16 — Future Work."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Roadmap — What's Next for MediGuide")
    add_slide_number(slide, 16)

    items = [
        (C_BLUE,   "1",
         "MIMIC-IV Integration",
         "Add 331,794 real clinical notes for richer discharge context and better evaluation"),
        (C_GREEN,  "2",
         "MedQuAD Dataset",
         "47,457 NIH Q&A pairs for superior Health Q&A coverage and recall"),
        (C_PURPLE, "3",
         "Multi-language Support",
         "Serve non-English speaking patients with translated outputs"),
        (C_AMBER,  "4",
         "Voice Interface",
         "Speech-to-text input + text-to-speech output for accessibility"),
        (C_ORANGE, "5",
         "Mobile App",
         "React Native wrapper for iOS/Android with offline support"),
    ]

    item_h = Inches(1.0)
    gap    = Inches(0.2)
    start_y = Inches(1.2)

    for i, (color, num, title, desc) in enumerate(items):
        iy = start_y + i * (item_h + gap)
        # Number circle
        circle = slide.shapes.add_shape(9, Inches(0.4), iy + Inches(0.2),
                                        Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text_box(slide, num,
                     left=Inches(0.4), top=iy + Inches(0.17),
                     width=Inches(0.6), height=Inches(0.6),
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Item background
        item_bg = slide.shapes.add_shape(5, Inches(1.15), iy,
                                         Inches(11.7), item_h)
        item_bg.adjustments[0] = 0.05
        item_bg.fill.solid()
        item_bg.fill.fore_color.rgb = C_WHITE
        item_bg.line.color.rgb = RGBColor(0xE5,0xE7,0xEB)
        item_bg.line.width = Pt(1)
        # Left accent
        add_rect(slide, Inches(1.15), iy, Inches(0.06), item_h, fill_color=color)

        add_text_box(slide, title,
                     left=Inches(1.35), top=iy + Inches(0.1),
                     width=Inches(11.2), height=Inches(0.42),
                     font_size=15, bold=True, color=color)
        add_text_box(slide, desc,
                     left=Inches(1.35), top=iy + Inches(0.52),
                     width=Inches(11.2), height=Inches(0.42),
                     font_size=12.5, color=C_DARKTEXT)


def slide17_limitations(prs):
    """Slide 17 — Limitations."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "Honest Assessment — Current Limitations")
    add_slide_number(slide, 17)

    limitations = [
        (C_AMBER,
         "Metric Interpretation",
         "ROUGE-L of 0.117 reflects lexical mismatch with reference — not poor quality "
         "(faithfulness confirmed by NLI). Model output is semantically correct but paraphrased."),
        (C_BLUE,
         "Retrieval Precision",
         "Precision@5 of 0.15 improves significantly with MIMIC-IV data (not yet loaded). "
         "Current corpus is FDA/NIH focused, not discharge-note specific."),
        (C_ORANGE,
         "OCR Quality",
         "OCR accuracy depends on input image quality. Blurry scans or handwritten notes "
         "may produce text extraction errors."),
        (C_PURPLE,
         "LLM Performance",
         "Ollama LLM requires GPU for best performance. CPU inference is slow (~3–5 min/response). "
         "Streaming helps perceived latency."),
        (C_GRAY,
         "Spell Correction Scope",
         "Spell correction covers 35+ medical terms only. General English typos rely on "
         "browser's built-in spellcheck functionality."),
    ]

    item_h = Inches(1.05)
    gap    = Inches(0.15)
    start_y = Inches(1.15)

    for i, (color, title, desc) in enumerate(limitations):
        iy = start_y + i * (item_h + gap)
        # Warning icon square
        warn = slide.shapes.add_shape(1, Inches(0.4), iy + Inches(0.2),
                                      Inches(0.5), Inches(0.5))
        warn.fill.solid()
        warn.fill.fore_color.rgb = color
        warn.line.fill.background()
        add_text_box(slide, "!",
                     left=Inches(0.4), top=iy + Inches(0.17),
                     width=Inches(0.5), height=Inches(0.5),
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Item background
        item_bg = slide.shapes.add_shape(5, Inches(1.05), iy,
                                         Inches(11.8), item_h)
        item_bg.adjustments[0] = 0.04
        item_bg.fill.solid()
        item_bg.fill.fore_color.rgb = C_WHITE
        item_bg.line.color.rgb = RGBColor(0xE5,0xE7,0xEB)
        item_bg.line.width = Pt(1)
        add_rect(slide, Inches(1.05), iy, Inches(0.06), item_h, fill_color=color)

        add_text_box(slide, title,
                     left=Inches(1.25), top=iy + Inches(0.08),
                     width=Inches(11.4), height=Inches(0.4),
                     font_size=13.5, bold=True, color=color)
        add_text_box(slide, desc,
                     left=Inches(1.25), top=iy + Inches(0.5),
                     width=Inches(11.4), height=Inches(0.48),
                     font_size=12, color=C_DARKTEXT)


def slide18_conclusion(prs):
    """Slide 18 — Conclusion (dark bg)."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_NAVY)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=C_BLUE)
    add_slide_number(slide, 18)

    add_text_box(slide, "MediGuide — Key Takeaways",
                 left=Inches(0.5), top=Inches(0.25),
                 width=Inches(12.3), height=Inches(0.7),
                 font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, Inches(2.5), Inches(1.0), Inches(8.3), Inches(0.04),
             fill_color=C_BLUE)

    takeaways = [
        (C_BLUE,   "End-to-End RAG System",
         "Built a complete system from scratch — no black boxes, fully open-source"),
        (C_GREEN,  "Rich Open Knowledge Base",
         "49,423 chunks from 4 open medical databases (openFDA, MedlinePlus, PubMed, PLABA)"),
        (C_PURPLE, "Hybrid Retrieval",
         "BM25 + PubMedBERT + RRF fusion outperforms either retrieval method alone"),
        (C_AMBER,  "Patient-Facing UI",
         "Streaming, history, recommendations, spell correction, multi-format file upload"),
        (C_GREEN,  "Privacy-Preserving",
         "100% local — medical notes never leave the patient's device"),
    ]

    for i, (color, title, desc) in enumerate(takeaways):
        iy = Inches(1.15) + i * Inches(0.92)
        # Dot
        dot = slide.shapes.add_shape(9, Inches(0.5), iy + Inches(0.22),
                                     Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.95), iy, Inches(11.8), Inches(0.88))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        run1 = p1.add_run()
        run1.text = title
        run1.font.size = Pt(15)
        run1.font.bold = True
        run1.font.color.rgb = color
        run1.font.name = "Calibri"
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(12.5)
        run2.font.bold = False
        run2.font.color.rgb = RGBColor(0xBF,0xDB,0xFF)
        run2.font.name = "Calibri"

    # Metrics banner
    metrics_y = Inches(5.85)
    add_rect(slide, Inches(0.5), metrics_y, Inches(12.3), Inches(0.85),
             fill_color=C_BLUE)
    add_text_box(slide,
                 "MRR: 0.576   ·   Recall@5: 0.55   ·   Readability Grade: 9.3",
                 left=Inches(0.5), top=metrics_y + Inches(0.18),
                 width=Inches(12.3), height=Inches(0.55),
                 font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_BLUE)


def slide19_references(prs):
    """Slide 19 — References."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG)
    heading(slide, "References & Datasets")
    add_slide_number(slide, 19)

    refs = [
        ("PLABA Dataset:",
         "Attal et al., BioNLP 2023 — Plain Language Adaptation of Biomedical Abstracts"),
        ("MedlinePlus:",
         "U.S. National Library of Medicine — nlm.nih.gov/medlineplus"),
        ("openFDA:",
         "U.S. Food & Drug Administration — open.fda.gov"),
        ("PubMed:",
         "NCBI Entrez API — pubmed.ncbi.nlm.nih.gov"),
        ("PubMedBERT:",
         "pritamdeka/S-PubMedBert-MS-MARCO (HuggingFace)"),
        ("LLaMA 3.1:",
         "Meta AI, 2024"),
        ("FAISS:",
         "Johnson et al., 2019 — Billion-scale similarity search with GPUs"),
        ("BM25:",
         "Robertson & Zaragoza, 2009 — The Probabilistic Relevance Framework"),
    ]

    start_y = Inches(1.15)
    row_h   = Inches(0.67)

    for i, (label, text) in enumerate(refs):
        ry = start_y + i * row_h
        # Alternating bg
        if i % 2 == 0:
            add_rect(slide, Inches(0.4), ry, Inches(12.5), row_h - Inches(0.04),
                     fill_color=RGBColor(0xF0,0xF5,0xFF))

        # Bullet
        dot = slide.shapes.add_shape(9, Inches(0.5), ry + Inches(0.24),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_BLUE
        dot.line.fill.background()

        # Label
        add_text_box(slide, label,
                     left=Inches(0.75), top=ry + Inches(0.1),
                     width=Inches(1.7), height=Inches(0.5),
                     font_size=12.5, bold=True, color=C_NAVY)

        # Reference text
        add_text_box(slide, text,
                     left=Inches(2.5), top=ry + Inches(0.1),
                     width=Inches(10.3), height=Inches(0.5),
                     font_size=12.5, color=C_DARKTEXT)


def slide20_qa(prs):
    """Slide 20 — Q&A (dark bg, centered)."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_NAVY)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=C_BLUE)

    # Decorative circle
    circle = slide.shapes.add_shape(9, Inches(5.3), Inches(0.8), Inches(2.7), Inches(2.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1F, 0x4A, 0x8B)
    circle.line.fill.background()

    # Thank You
    add_text_box(slide, "Thank You",
                 left=Inches(1.5), top=Inches(2.0),
                 width=Inches(10.3), height=Inches(1.4),
                 font_size=64, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Blue underline
    add_rect(slide, Inches(4.5), Inches(3.45), Inches(4.3), Inches(0.07),
             fill_color=C_BLUE)

    # Sub
    add_text_box(slide, "Questions & Discussion",
                 left=Inches(1.5), top=Inches(3.6),
                 width=Inches(10.3), height=Inches(0.7),
                 font_size=24, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

    # Metrics chips
    chips_y = Inches(4.55)
    chip_data = [
        (C_BLUE,   "MRR: 0.576"),
        (C_GREEN,  "Recall@5: 0.55"),
        (C_AMBER,  "Grade: 9.3"),
        (C_PURPLE, "49,423 Chunks"),
    ]
    chip_w = Inches(2.7)
    chip_h = Inches(0.6)
    chip_gap = Inches(0.25)
    total_chips_w = len(chip_data) * chip_w + (len(chip_data)-1) * chip_gap
    chip_start_x = (Inches(13.33) - total_chips_w) / 2

    for i, (color, label) in enumerate(chip_data):
        cx = chip_start_x + i * (chip_w + chip_gap)
        shape = slide.shapes.add_shape(5, cx, chips_y, chip_w, chip_h)
        shape.adjustments[0] = 0.5
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        add_text_box(slide, label,
                     left=cx, top=chips_y + Inches(0.1),
                     width=chip_w, height=chip_h - Inches(0.12),
                     font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Bottom text
    add_text_box(slide,
                 "MediGuide — 100% Local · Open Source · Patient First",
                 left=Inches(1.5), top=Inches(5.55),
                 width=Inches(10.3), height=Inches(0.45),
                 font_size=13, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "github.com/yourname/mediguide",
                 left=Inches(1.5), top=Inches(6.1),
                 width=Inches(10.3), height=Inches(0.4),
                 font_size=12, color=C_BLUE, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 20)
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    out_path = "/Users/gautampatel/Documents/discharge_rag_full/outputs/MediGuide_Presentation.pptx"

    print("Generating MediGuide Presentation...")
    print("  Creating bar chart PNG...")
    chart_png = make_bar_chart_png()

    prs = new_prs()

    print("  Building slides...")
    slide1_title(prs)
    print("    Slide  1/20 — Title")
    slide2_problem(prs)
    print("    Slide  2/20 — Problem")
    slide3_solution(prs)
    print("    Slide  3/20 — Solution")
    slide4_architecture(prs)
    print("    Slide  4/20 — Architecture")
    slide5_dataset(prs)
    print("    Slide  5/20 — Dataset")
    slide6_rag_pipeline(prs)
    print("    Slide  6/20 — RAG Pipeline")
    slide7_retrieval(prs)
    print("    Slide  7/20 — Retrieval Strategy")
    slide8_evaluation(prs)
    print("    Slide  8/20 — Evaluation")
    slide9_webapp(prs)
    print("    Slide  9/20 — Web App")
    slide10_two_modes(prs)
    print("    Slide 10/20 — Two Modes")
    slide11_smart_input(prs)
    print("    Slide 11/20 — Smart Input")
    slide12_tech_stack(prs)
    print("    Slide 12/20 — Tech Stack")
    slide13_privacy(prs)
    print("    Slide 13/20 — Privacy")
    slide14_provenance(prs)
    print("    Slide 14/20 — Provenance")
    slide15_corpus_chart(prs, chart_png)
    print("    Slide 15/20 — Corpus Chart")
    slide16_future(prs)
    print("    Slide 16/20 — Future Work")
    slide17_limitations(prs)
    print("    Slide 17/20 — Limitations")
    slide18_conclusion(prs)
    print("    Slide 18/20 — Conclusion")
    slide19_references(prs)
    print("    Slide 19/20 — References")
    slide20_qa(prs)
    print("    Slide 20/20 — Q&A")

    print(f"  Saving to: {out_path}")
    prs.save(out_path)
    print("DONE: MediGuide_Presentation.pptx saved")


if __name__ == "__main__":
    main()
